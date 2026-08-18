"""PowerPoint rendering.

A deck is not a document: a slide that needs scrolling is a failed slide. So
each section becomes a narrative-and-metrics slide, and every table is split
across as many slides as it needs, with the header repeated and the slide title
numbered ("... (2 of 3)").

Built on the blank layout with explicit text boxes rather than placeholder
layouts, because placeholder sets differ between templates and would make the
output depend on python-pptx's bundled theme.
"""

from __future__ import annotations

import io
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.presentation import Presentation as Deck
from pptx.util import Inches, Length, Pt

from app.schemas.report import ReportData, ReportSection, ReportTable
from app.services.report_builder import cell_text

CONTENT_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

#: 16:9, the modern default.
_SLIDE_WIDTH = Inches(13.333)
_SLIDE_HEIGHT = Inches(7.5)
_MARGIN = Inches(0.6)
_CONTENT_WIDTH = _SLIDE_WIDTH - (2 * _MARGIN)
_BODY_TOP = Inches(1.5)

_ACCENT = RGBColor(0x1F, 0x38, 0x64)
_MUTED = RGBColor(0x5A, 0x64, 0x72)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_BAND = RGBColor(0xF2, 0xF4, 0xF8)

#: Table rows per slide, and bullets per slide.
_ROWS_PER_SLIDE = 10
_BULLETS_PER_SLIDE = 8
#: Bullets longer than this are trimmed; a slide is not a paragraph.
_MAX_BULLET_CHARS = 220


def _blank(presentation: Deck) -> Any:
    return presentation.slides.add_slide(presentation.slide_layouts[6])


def _textbox(
    slide: Any,
    text: str,
    *,
    top: Length,
    size: float,
    bold: bool = False,
    color: RGBColor = _ACCENT,
    height: Length | None = None,
) -> Any:
    box = slide.shapes.add_textbox(
        _MARGIN, top, _CONTENT_WIDTH, height or Inches(0.6)
    )
    frame = box.text_frame
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _title_slide(presentation: Deck, data: ReportData) -> None:
    slide = _blank(presentation)
    _textbox(slide, data.title, top=Inches(2.2), size=36, bold=True, height=Inches(1.4))
    if data.subtitle:
        _textbox(slide, data.subtitle, top=Inches(3.5), size=16, color=_MUTED, height=Inches(1.0))

    facts = " • ".join(
        [
            data.dataset_name,
            data.version_label,
            f"{data.row_count:,} rows × {data.column_count:,} columns",
            data.generated_at.strftime("%Y-%m-%d"),
            data.generated_by,
        ]
    )
    _textbox(slide, facts, top=Inches(4.6), size=12, color=_MUTED, height=Inches(0.8))


def _agenda_slide(presentation: Deck, data: ReportData) -> None:
    slide = _blank(presentation)
    _textbox(slide, "Contents", top=_MARGIN, size=28, bold=True)
    box = slide.shapes.add_textbox(_MARGIN, _BODY_TOP, _CONTENT_WIDTH, Inches(5.0))
    frame = box.text_frame
    frame.word_wrap = True
    for index, section in enumerate(data.sections, start=1):
        paragraph = frame.paragraphs[0] if index == 1 else frame.add_paragraph()
        run = paragraph.add_run()
        run.text = f"{index}.  {section.title}"
        run.font.size = Pt(16)
        run.font.color.rgb = _ACCENT
        paragraph.space_after = Pt(8)


def _bullet_slides(presentation: Deck, section: ReportSection) -> None:
    """Narrative, metrics and bullets - as many slides as the content needs."""
    lines: list[tuple[str, bool]] = [(text, False) for text in section.narrative]
    if section.unavailable_reason:
        lines.append((section.unavailable_reason, False))
    lines.extend(
        (f"{metric.label}: {metric.value}" + (f" ({metric.detail})" if metric.detail else ""), True)
        for metric in section.metrics
    )
    lines.extend((bullet, True) for bullet in section.bullets)

    if not lines:
        return

    chunks = [
        lines[index : index + _BULLETS_PER_SLIDE]
        for index in range(0, len(lines), _BULLETS_PER_SLIDE)
    ]
    for position, chunk in enumerate(chunks, start=1):
        slide = _blank(presentation)
        suffix = f" ({position} of {len(chunks)})" if len(chunks) > 1 else ""
        _textbox(slide, f"{section.title}{suffix}", top=_MARGIN, size=26, bold=True)

        box = slide.shapes.add_textbox(_MARGIN, _BODY_TOP, _CONTENT_WIDTH, Inches(5.2))
        frame = box.text_frame
        frame.word_wrap = True
        for index, (text, is_bullet) in enumerate(chunk):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            run = paragraph.add_run()
            trimmed = text if len(text) <= _MAX_BULLET_CHARS else f"{text[:_MAX_BULLET_CHARS]}…"
            run.text = f"•  {trimmed}" if is_bullet else trimmed
            run.font.size = Pt(14 if is_bullet else 13)
            run.font.color.rgb = _ACCENT if is_bullet else _MUTED
            paragraph.space_after = Pt(8)


def _table_slides(presentation: Deck, section: ReportSection, table: ReportTable) -> None:
    if not table.columns or not table.rows:
        return

    chunks = [
        table.rows[index : index + _ROWS_PER_SLIDE]
        for index in range(0, len(table.rows), _ROWS_PER_SLIDE)
    ]
    heading = f"{section.title} — {table.title}" if table.title else section.title
    # Narrower columns need smaller type to stay inside the cell.
    font_size = 11 if len(table.columns) <= 5 else 9 if len(table.columns) <= 8 else 7

    for position, chunk in enumerate(chunks, start=1):
        slide = _blank(presentation)
        suffix = f" ({position} of {len(chunks)})" if len(chunks) > 1 else ""
        _textbox(slide, f"{heading}{suffix}", top=_MARGIN, size=24, bold=True)

        shape = slide.shapes.add_table(
            len(chunk) + 1,
            len(table.columns),
            _MARGIN,
            _BODY_TOP,
            _CONTENT_WIDTH,
            Inches(0.4) * (len(chunk) + 1),
        )
        grid = shape.table

        for index, column in enumerate(table.columns):
            cell = grid.cell(0, index)
            cell.text = str(column)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.runs[0].font.size = Pt(font_size)
            paragraph.runs[0].font.bold = True
            paragraph.runs[0].font.color.rgb = _WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = _ACCENT

        for row_index, row in enumerate(chunk, start=1):
            values = list(row)[: len(table.columns)]
            values += [None] * (len(table.columns) - len(values))
            for column_index, value in enumerate(values):
                cell = grid.cell(row_index, column_index)
                cell.text = cell_text(value)
                paragraph = cell.text_frame.paragraphs[0]
                if paragraph.runs:
                    paragraph.runs[0].font.size = Pt(font_size)
                cell.fill.solid()
                cell.fill.fore_color.rgb = _WHITE if row_index % 2 else _BAND

        if table.note and position == len(chunks):
            _textbox(
                slide,
                table.note,
                top=Inches(6.7),
                size=10,
                color=_MUTED,
                height=Inches(0.4),
            )


def _skipped_slide(presentation: Deck, data: ReportData) -> None:
    if not data.skipped:
        return
    slide = _blank(presentation)
    _textbox(slide, "Sections not included", top=_MARGIN, size=26, bold=True)
    box = slide.shapes.add_textbox(_MARGIN, _BODY_TOP, _CONTENT_WIDTH, Inches(5.2))
    frame = box.text_frame
    frame.word_wrap = True
    for index, item in enumerate(data.skipped[:10]):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        run = paragraph.add_run()
        run.text = f"•  {item['section']}: {item['reason']}"
        run.font.size = Pt(12)
        run.font.color.rgb = _MUTED
        paragraph.space_after = Pt(6)


def render(data: ReportData) -> bytes:
    """Render the report as a 16:9 slide deck."""
    presentation = Presentation()
    presentation.slide_width = _SLIDE_WIDTH
    presentation.slide_height = _SLIDE_HEIGHT

    _title_slide(presentation, data)
    if data.sections:
        _agenda_slide(presentation, data)

    for section in data.sections:
        _bullet_slides(presentation, section)
        for table in section.tables:
            _table_slides(presentation, section, table)

    _skipped_slide(presentation, data)

    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()
